"""Create PowerPoint tool."""
from typing import Any, Dict, List, Optional
from pathlib import Path
from .base import Tool, ToolResult


class CreatePPTXTool(Tool):
    """Tool for creating PowerPoint presentations."""

    def __init__(self):
        super().__init__(
            name="create_pptx",
            description="Create a PowerPoint presentation. Input is a JSON string with slides array containing 'title' and 'content' for each slide.",
        )

    async def execute(self, json_str: str, **kwargs) -> ToolResult:
        try:
            import json
            import os
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN

            data = json.loads(json_str)
            slides = data.get("slides", [])
            
            if not slides:
                return ToolResult(
                    success=False, 
                    result=None, 
                    error="No slides provided"
                )

            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
            
            title_slide_layout = prs.slide_layouts[6]
            title_content_layout = prs.slide_layouts[1]

            for slide_data in slides:
                slide = prs.slides.add_slide(title_content_layout)
                
                title = slide.shapes.title
                title.text = slide_data.get("title", "")
                title.text_frame.paragraphs[0].font.size = Pt(32)
                title.text_frame.paragraphs[0].font.bold = True
                
                content_box = slide.placeholders[1]
                content = slide_data.get("content", "")
                
                if isinstance(content, list):
                    content_text = "\n".join(f"• {item}" for item in content)
                else:
                    content_text = content
                    
                content_box.text = content_text
                content_box.text_frame.paragraphs[0].font.size = Pt(18)

            output_path = data.get("output", "presentation.pptx")
            prs.save(output_path)
            
            return ToolResult(
                success=True, 
                result=f"Created {output_path} with {len(slides)} slides"
            )

        except ImportError:
            return ToolResult(
                success=False, 
                result=None, 
                error="python-pptx not installed. Run: pip install python-pptx"
            )
        except Exception as e:
            return ToolResult(
                success=False, 
                result=None, 
                error=str(e)
            )

    def _get_parameters_schema(self) -> Dict[str, Any]:
        return {
            "type": "object",
            "properties": {
                "json_str": {
                    "type": "string",
                    "description": 'JSON string with slides: {"slides": [{"title": "Slide 1", "content": "..."}, ...], "output": "file.pptx"}',
                }
            },
            "required": ["json_str"],
        }